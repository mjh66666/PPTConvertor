# -*- coding: utf-8 -*-


import os
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
import win32com.client
from PySide6.QtWidgets import (QApplication, QMainWindow, QVBoxLayout, QHBoxLayout,
                               QWidget, QLabel, QLineEdit, QPushButton, QCheckBox,
                               QRadioButton, QGroupBox, QFileDialog, QMessageBox)
from PySide6.QtCore import Qt, QSize
from PySide6.QtGui import QIcon

# ==================== 常量定义 ====================
APP_NAME = "PPT转换工具"
VERSION = "1.0.0"
AUTHOR = "莫炯豪"
COMPANY = "萝卜丁团队"
YEAR = "2025"


# PowerPoint保存格式常量
PP_SAVE_AS_PNG = 18  # PNG格式
PP_SAVE_AS_JPG = 17  # JPG格式

import sys
import os
from PySide6.QtGui import QIcon


def resource_path(relative_path):
    """获取资源文件路径（兼容PyInstaller打包后的路径）"""
    try:
        # PyInstaller打包后的临时路径
        base_path = sys._MEIPASS
    except AttributeError:
        base_path = os.path.abspath(".")
    return os.path.join(base_path, relative_path)


# ==================== 主窗口类 ====================
class PPTConverterUI(QMainWindow):
    """主应用程序窗口，提供PPT转换功能"""

    def __init__(self):
        super().__init__()

        # 窗口基本设置
        self.setWindowTitle(f"{APP_NAME} v{VERSION}")
        self.setMinimumSize(600, 450)

        # 设置窗口图标
        self.setup_icons()

        # 初始化UI
        self.init_ui()

        # 连接信号槽
        self.connect_signals()

        # 存储文件大小信息
        self.original_size = 0
        self.final_size = 0

    def setup_icons(self):
        self.setWindowIcon(QIcon(resource_path("icons/lbd128.ico")))
    def init_ui(self):
        """初始化用户界面"""
        # 主部件和布局
        main_widget = QWidget()
        main_layout = QVBoxLayout(main_widget)
        main_layout.setContentsMargins(15, 15, 15, 15)
        main_layout.setSpacing(15)
        self.setWindowIcon(QIcon("../icons/lbd128.ico"))
        # 1. 输入PPT文件选择
        input_group = QGroupBox("1. 选择输入PPT文件")
        input_layout = QHBoxLayout()
        self.input_path = QLineEdit()
        self.input_path.setPlaceholderText("请选择PPT文件...")
        self.input_path.setReadOnly(True)
        self.input_btn = QPushButton("浏览...")
        self.input_btn.setFixedWidth(80)
        input_layout.addWidget(self.input_path)
        input_layout.addWidget(self.input_btn)
        input_group.setLayout(input_layout)

        # 2. 图片格式选择
        format_group = QGroupBox("2. 选择图片格式")
        format_layout = QHBoxLayout()
        self.png_radio = QRadioButton("PNG (高质量)")
        self.jpg_radio = QRadioButton("JPG (较小文件)")
        self.png_radio.setChecked(True)
        format_layout.addWidget(self.png_radio)
        format_layout.addWidget(self.jpg_radio)
        format_group.setLayout(format_layout)

        # 3. 临时图片文件夹选择
        temp_group = QGroupBox("3. 选择临时图片保存位置")
        temp_layout = QHBoxLayout()
        self.temp_path = QLineEdit()
        self.temp_path.setPlaceholderText("请选择临时图片保存文件夹...")
        self.temp_path.setReadOnly(True)
        self.temp_btn = QPushButton("浏览...")
        self.temp_btn.setFixedWidth(80)
        temp_layout.addWidget(self.temp_path)
        temp_layout.addWidget(self.temp_btn)
        temp_group.setLayout(temp_layout)

        # 4. 输出PPT文件选择
        output_group = QGroupBox("4. 选择输出PPT位置和名称")
        output_layout = QHBoxLayout()
        self.output_path = QLineEdit()
        self.output_path.setPlaceholderText("请选择输出PPT文件...")
        self.output_path.setReadOnly(True)
        self.output_btn = QPushButton("浏览...")
        self.output_btn.setFixedWidth(80)
        output_layout.addWidget(self.output_path)
        output_layout.addWidget(self.output_btn)
        output_group.setLayout(output_layout)

        # 5. 选项设置
        options_group = QGroupBox("选项设置")
        options_layout = QVBoxLayout()

        # 删除临时文件选项
        self.delete_temp_check = QCheckBox("转换完成后删除临时图片文件")
        self.delete_temp_check.setChecked(True)

        # 打开输出文件夹选项
        self.open_output_check = QCheckBox("转换完成后打开输出文件夹")
        self.open_output_check.setChecked(True)

        options_layout.addWidget(self.delete_temp_check)
        options_layout.addWidget(self.open_output_check)
        options_group.setLayout(options_layout)

        # 6. 转换按钮
        self.convert_btn = QPushButton("开始转换")
        self.convert_btn.setFixedHeight(40)
        self.convert_btn.setStyleSheet("""
            QPushButton {
                background-color: #4CAF50;
                color: white;
                font-weight: bold;
                border-radius: 5px;
            }
            QPushButton:hover {
                background-color: #45a049;
            }
            QPushButton:disabled {
                background-color: #cccccc;
            }
        """)

        # 7. 底部信息栏
        footer_layout = QHBoxLayout()

        # 版本信息
        version_label = QLabel(f"v{VERSION} © {YEAR} {COMPANY}")
        version_label.setStyleSheet("color: #666666;")

        # 警告信息
        warning_label = QLabel("注意: 转换过程中PPT会被临时打开，请勿关闭！")
        warning_label.setStyleSheet("color: red; font-weight: bold;")

        footer_layout.addWidget(version_label)
        footer_layout.addStretch()
        footer_layout.addWidget(warning_label)

        # 添加到主布局
        main_layout.addWidget(input_group)
        main_layout.addWidget(format_group)
        main_layout.addWidget(temp_group)
        main_layout.addWidget(output_group)
        main_layout.addWidget(options_group)
        main_layout.addWidget(self.convert_btn)
        main_layout.addLayout(footer_layout)

        self.setCentralWidget(main_widget)

    def connect_signals(self):
        """连接信号与槽"""
        self.input_btn.clicked.connect(self.select_input_file)
        self.temp_btn.clicked.connect(self.select_temp_folder)
        self.output_btn.clicked.connect(self.select_output_file)
        self.convert_btn.clicked.connect(self.start_conversion)

    # ==================== 业务逻辑方法 ====================
    def select_input_file(self):
        """选择输入PPT文件"""
        file_path, _ = QFileDialog.getOpenFileName(
            self, "选择PPT文件", "",
            "PowerPoint文件 (*.pptx *.ppt);;所有文件 (*.*)"
        )
        if file_path:
            self.input_path.setText(file_path)
            # 获取原始文件大小
            self.original_size = os.path.getsize(file_path)

    def select_temp_folder(self):
        """选择临时图片保存文件夹"""
        folder_path = QFileDialog.getExistingDirectory(
            self, "选择临时图片保存文件夹"
        )
        if folder_path:
            self.temp_path.setText(folder_path)

    def select_output_file(self):
        """选择输出PPT文件位置"""
        file_path, _ = QFileDialog.getSaveFileName(
            self, "保存PPT文件", "",
            "PowerPoint文件 (*.pptx);;所有文件 (*.*)"
        )
        if file_path:
            # 确保扩展名正确
            if not file_path.lower().endswith('.pptx'):
                file_path += '.pptx'
            self.output_path.setText(file_path)

    def start_conversion(self):
        """开始转换流程"""
        # 验证输入
        if not self.validate_inputs():
            return

        # 显示警告信息
        if not self.show_warning_message():
            return

        # 获取参数
        input_ppt = self.input_path.text()
        temp_dir = self.temp_path.text()
        output_ppt = self.output_path.text()
        img_format = PP_SAVE_AS_PNG if self.png_radio.isChecked() else PP_SAVE_AS_JPG
        delete_temp = self.delete_temp_check.isChecked()
        open_output = self.open_output_check.isChecked()

        # 禁用按钮防止重复点击
        self.convert_btn.setEnabled(False)

        try:
            # 显示开始消息
            QMessageBox.information(
                self, "信息", "即将开始转换过程，请耐心等待..."
            )

            # 1. PPT转图片
            if not self.ppt_to_images(input_ppt, temp_dir, img_format):
                QMessageBox.critical(
                    self, "错误", "PPT转换为图片失败！"
                )
                return

            # 2. 图片转PPT
            if not self.images_to_ppt(temp_dir, output_ppt):
                QMessageBox.critical(
                    self, "错误", "图片转换为PPT失败！"
                )
                return

            # 3. 获取最终文件大小
            self.final_size = os.path.getsize(output_ppt)

            # 4. 可选删除临时文件
            if delete_temp:
                self.delete_temp_files(temp_dir)

            # 5. 显示完成信息
            self.show_completion_message()

            # 6. 可选打开输出文件夹
            if open_output:
                output_dir = os.path.dirname(output_ppt)
                os.startfile(output_dir)

        finally:
            # 重新启用按钮
            self.convert_btn.setEnabled(True)

    def validate_inputs(self):
        """验证输入是否完整"""
        errors = []

        if not self.input_path.text():
            errors.append("请选择输入PPT文件！")

        if not self.temp_path.text():
            errors.append("请选择临时图片保存文件夹！")

        if not self.output_path.text():
            errors.append("请选择输出PPT文件！")

        if errors:
            QMessageBox.warning(
                self, "输入不完整", "\n".join(errors)
            )
            return False

        return True

    def show_warning_message(self):
        """显示转换警告信息"""
        msg = QMessageBox(self)
        msg.setIcon(QMessageBox.Warning)
        msg.setWindowTitle("重要提示")
        msg.setText("转换过程中PPT会被临时打开")
        msg.setInformativeText(
            "请确保:\n"
            "1. 已关闭所有PowerPoint窗口\n"
            "2. 不要关闭转换过程中弹出的PowerPoint窗口\n"
            "3. 转换期间不要操作电脑"
        )
        msg.setStandardButtons(QMessageBox.Ok | QMessageBox.Cancel)
        msg.setDefaultButton(QMessageBox.Ok)

        return msg.exec() == QMessageBox.Ok

    def ppt_to_images(self, ppt_path, output_dir, img_format):
        """将PPT转换为图片序列"""
        try:
            ppt_path = str(Path(ppt_path).resolve())
            output_dir = str(Path(output_dir).resolve())

            if not os.path.exists(ppt_path):
                raise FileNotFoundError(f"PPT文件不存在: {ppt_path}")

            os.makedirs(output_dir, exist_ok=True)

            # 启动PowerPoint
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            powerpoint.Visible = 1  # 设为可见便于调试

            # 打开演示文稿
            presentation = powerpoint.Presentations.Open(ppt_path)

            # 导出幻灯片为图片
            total_slides = presentation.Slides.Count
            for i, slide in enumerate(presentation.Slides, start=1):
                output_file = os.path.join(output_dir, f"Slide_{i}.png")
                slide.Export(
                    output_file,
                    "PNG" if img_format == PP_SAVE_AS_PNG else "JPG"
                )

                # 更新进度（可选）
                print(f"已导出 {i}/{total_slides}")

            presentation.Close()
            powerpoint.Quit()
            return True

        except Exception as e:
            print(f"转换失败: {str(e)}")
            if 'presentation' in locals():
                presentation.Close()
            if 'powerpoint' in locals():
                powerpoint.Quit()
            return False

    def images_to_ppt(self, image_folder, output_ppt_path):
        """将图片序列转换为PPT"""
        try:
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(5.625)

            # 获取文件夹中所有Slide_*.png文件并按数字排序
            slide_files = sorted(
                [f for f in os.listdir(image_folder)
                 if f.startswith("Slide_") and f.endswith((".png", ".jpg"))],
                key=lambda x: int(x.split("_")[1].split(".")[0])
            )

            total_slides = len(slide_files)
            for i, slide_file in enumerate(slide_files, start=1):
                img_path = os.path.join(image_folder, slide_file)

                # 添加空白幻灯片
                slide = prs.slides.add_slide(prs.slide_layouts[6])

                # 插入图片并居中显示
                pic = slide.shapes.add_picture(
                    img_path, 0, 0,
                    width=prs.slide_width,
                    height=prs.slide_height
                )

                # 锁定纵横比并居中
                pic.lock_aspect_ratio = True
                pic.left = int((prs.slide_width - pic.width) / 2)
                pic.top = int((prs.slide_height - pic.height) / 2)

                # 更新进度（可选）
                print(f"已处理 {i}/{total_slides}")

            prs.save(output_ppt_path)
            return True

        except Exception as e:
            print(f"创建PPT失败: {str(e)}")
            return False

    def delete_temp_files(self, temp_dir):
        """删除临时图片文件"""
        try:
            deleted_files = 0
            for file in os.listdir(temp_dir):
                if file.startswith("Slide_") and file.endswith((".png", ".jpg")):
                    os.remove(os.path.join(temp_dir, file))
                    deleted_files += 1
            print(f"已删除 {deleted_files} 个临时文件")
        except Exception as e:
            print(f"删除临时文件失败: {str(e)}")

    def show_completion_message(self):
        """显示转换完成信息"""

        def format_size(size):
            """格式化文件大小"""
            for unit in ['B', 'KB', 'MB', 'GB']:
                if size < 1024.0:
                    return f"{size:.2f} {unit}"
                size /= 1024.0
            return f"{size:.2f} TB"

        original_size_fmt = format_size(self.original_size)
        final_size_fmt = format_size(self.final_size)

        msg = QMessageBox(self)
        msg.setIcon(QMessageBox.Information)
        msg.setWindowTitle("转换完成")
        msg.setText("PPT转换已完成！")
        msg.setInformativeText(
            f"原始文件大小: {original_size_fmt}\n"
            f"最终文件大小: {final_size_fmt}\n\n"
            f"文件已保存到:\n{self.output_path.text()}"
        )
        msg.exec()


# ==================== 应用程序入口 ====================
def main():
    """应用程序主入口"""
    try:
        # 创建应用实例
        app = QApplication(sys.argv)
        app.setApplicationName(APP_NAME)
        app.setApplicationVersion(VERSION)
        app.setOrganizationName(COMPANY)

        # 创建并显示主窗口
        converter = PPTConverterUI()
        converter.show()

        # 运行应用
        sys.exit(app.exec())

    except Exception as e:
        print(f"应用程序错误: {str(e)}")
        QMessageBox.critical(
            None, "致命错误",
            f"应用程序遇到错误:\n{str(e)}\n\n程序将退出。"
        )
        sys.exit(1)


def handle_exceptions(exc_type, exc_value, exc_traceback):
    """捕获所有未处理的异常"""
    with open("error.log", "a") as f:
        f.write(f"CRASH: {str(exc_value)}\n")


if __name__ == '__main__':
    sys.excepthook = handle_exceptions
    main()
